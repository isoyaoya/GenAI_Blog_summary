#!/usr/bin/env python3
"""
PowerPoint生成器模块
用于将AWS博客摘要转换为PowerPoint演示文稿

依赖项：
- python-pptx: PowerPoint文件生成
"""

from typing import Dict, List, Optional, Any
from dataclasses import dataclass
from datetime import datetime
import logging

# 可选导入python-pptx，如果不可用则优雅降级
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx库未安装，PowerPoint生成功能将不可用。请运行: pip install python-pptx")

@dataclass
class BlogSummaryData:
    """博客摘要数据结构"""
    title: str
    date: str
    url: str
    summary: str
    interest_score: int
    content_length: int = 0

@dataclass
class PresentationMetadata:
    """演示文稿元数据"""
    target_month: str
    generation_date: str
    total_articles: int
    interested_articles: int
    filtering_efficiency: float
    customer_info: Dict[str, str]

@dataclass
class TemplateConfig:
    """模板配置"""
    theme_name: str = "default"
    primary_color: str = "#1f4e79"
    secondary_color: str = "#70ad47"
    font_family: str = "Calibri"
    title_font_size: int = 16
    subtitle_font_size: int = 12
    content_font_size: int = 10
    url_font_size: int = 9

class PowerPointGenerator:
    """PowerPoint演示文稿生成器"""
    
    def __init__(self, template_config: Optional[TemplateConfig] = None):
        """
        初始化PowerPoint生成器
        
        Args:
            template_config: 可选的模板配置
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx库未安装，无法使用PowerPoint生成功能")
        
        self.template_config = template_config or TemplateConfig()
        self.presentation = None
        
    def create_presentation(self, summaries_data: List[BlogSummaryData], metadata: PresentationMetadata) -> bool:
        """
        创建完整的PowerPoint演示文稿
        
        Args:
            summaries_data: 博客摘要数据列表
            metadata: 演示文稿元数据
            
        Returns:
            bool: 创建是否成功
        """
        try:
            self.presentation = Presentation()
            
            # 设置幻灯片尺寸为16:9
            self.presentation.slide_width = Inches(13.33)  # 16:9比例的宽度
            self.presentation.slide_height = Inches(7.5)   # 16:9比例的高度
            
            # 添加标题幻灯片
            self._add_title_slide(metadata)
            
            # 添加摘要统计幻灯片
            self._add_statistics_slide(metadata)
            
            # 为每个博客摘要添加幻灯片
            for i, summary in enumerate(summaries_data, 1):
                self._add_summary_slide(summary, i)
            
            # 添加结论幻灯片
            self._add_conclusion_slide(metadata)
            
            return True
            
        except Exception as e:
            logging.error(f"创建PowerPoint演示文稿失败: {e}")
            return False
    
    def _add_title_slide(self, metadata: PresentationMetadata):
        """添加标题幻灯片"""
        slide_layout = self.presentation.slide_layouts[0]  # 标题幻灯片布局
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = f"AWS机器学习博客摘要报告"
        
        # 设置副标题
        subtitle = slide.placeholders[1]
        subtitle.text = f"{metadata.target_month} | 生成日期: {metadata.generation_date}"
        
        # 应用样式
        self._apply_title_styling(title)
        self._apply_subtitle_styling(subtitle)
    
    def _add_statistics_slide(self, metadata: PresentationMetadata):
        """添加统计信息幻灯片"""
        slide_layout = self.presentation.slide_layouts[1]  # 标题和内容布局
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = "📊 分析统计概览"
        
        # 添加统计内容
        content = slide.placeholders[1]
        stats_text = f"""📈 处理统计:
• 总处理文章数: {metadata.total_articles} 篇
• 客户感兴趣文章数: {metadata.interested_articles} 篇
• 筛选效率: {metadata.filtering_efficiency:.1f}%

🏢 客户信息概览:
• 公司名称: {metadata.customer_info.get('company', '合合信息')}
• 核心业务: OCR识别、文档智能化处理、AI智能终端
• 技术架构: AWS EKS + EC2, AWS Bedrock Claude模型
• 应用场景: 文档处理、图像识别、云资源管理

🎯 分析目标:
• 识别与客户业务相关的AWS技术趋势
• 发现可应用于现有产品的新技术
• 为技术决策提供数据支撑"""
        content.text = stats_text.strip()
        
        # 应用样式
        self._apply_title_styling(title)
        self._apply_content_styling(content)
    
    def _add_summary_slide(self, summary: BlogSummaryData, slide_number: int):
        """添加博客摘要幻灯片"""
        slide_layout = self.presentation.slide_layouts[1]  # 标题和内容布局
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = f"{slide_number}. {self._truncate_title(summary.title)}"
        
        # 添加摘要内容
        content = slide.placeholders[1]
        
        # 格式化摘要内容，分段显示
        formatted_summary = self._format_summary_for_slide(summary.summary)
        
        content_text = f"""📅 发布日期: {summary.date}

📝 内容摘要:
{formatted_summary}"""
        
        content.text = content_text.strip()
        
        # 应用样式
        self._apply_title_styling(title)
        self._apply_content_styling(content)
        
        # 添加原文链接作为单独的文本框
        self._add_url_textbox(slide, summary.url)
    
    def _add_conclusion_slide(self, metadata: PresentationMetadata):
        """添加结论幻灯片"""
        slide_layout = self.presentation.slide_layouts[1]  # 标题和内容布局
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = "💡 总结与建议"
        
        # 添加结论内容
        content = slide.placeholders[1]
        conclusion_text = f"""📋 分析总结:
本次分析共处理了 {metadata.total_articles} 篇AWS机器学习博客文章，筛选出 {metadata.interested_articles} 篇与客户业务高度相关的文章，筛选效率达到 {metadata.filtering_efficiency:.1f}%。

🔍 主要发现:
• AWS在OCR和文档处理领域持续创新
• Bedrock平台为AI应用提供更强大的基础能力
• 容器化和云原生技术成为AI部署的主流趋势
• 多模态AI模型在文档理解方面表现突出

📈 技术趋势:
• 大模型与传统OCR技术的深度融合
• 边缘计算在文档处理中的应用增加
• 自动化运维工具需求持续增长

🎯 行动建议:
• 深入研究相关文章中的技术方案
• 评估新技术在现有产品中的集成可能性
• 制定技术验证和实施路线图
• 持续跟踪AWS技术更新动态"""
        content.text = conclusion_text.strip()
        
        # 应用样式
        self._apply_title_styling(title)
        self._apply_content_styling(content)
    
    def _apply_title_styling(self, title_shape):
        """应用标题样式"""
        if title_shape.has_text_frame:
            text_frame = title_shape.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph.font.name = self.template_config.font_family
                paragraph.font.size = Pt(self.template_config.title_font_size)
                paragraph.font.bold = True
                # 解析颜色字符串
                color_hex = self.template_config.primary_color.lstrip('#')
                if len(color_hex) == 6:
                    r = int(color_hex[0:2], 16)
                    g = int(color_hex[2:4], 16)
                    b = int(color_hex[4:6], 16)
                    paragraph.font.color.rgb = RGBColor(r, g, b)
    
    def _apply_subtitle_styling(self, subtitle_shape):
        """应用副标题样式"""
        if subtitle_shape.has_text_frame:
            text_frame = subtitle_shape.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph.font.name = self.template_config.font_family
                paragraph.font.size = Pt(self.template_config.subtitle_font_size)
                # 解析颜色字符串
                color_hex = self.template_config.secondary_color.lstrip('#')
                if len(color_hex) == 6:
                    r = int(color_hex[0:2], 16)
                    g = int(color_hex[2:4], 16)
                    b = int(color_hex[4:6], 16)
                    paragraph.font.color.rgb = RGBColor(r, g, b)
    
    def _apply_content_styling(self, content_shape):
        """应用内容样式"""
        if content_shape.has_text_frame:
            text_frame = content_shape.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = None  # 禁用自动调整大小
            
            for paragraph in text_frame.paragraphs:
                paragraph.font.name = self.template_config.font_family
                paragraph.font.size = Pt(self.template_config.content_font_size)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)  # 黑色文本
                paragraph.line_spacing = 1.0  # 减小行间距
    
    def _truncate_title(self, title: str, max_length: int = 50) -> str:
        """截断标题以适应幻灯片"""
        if len(title) <= max_length:
            return title
        
        # 尝试在合适的位置截断（如空格、标点符号）
        truncated = title[:max_length-3]
        
        # 寻找最后一个空格或标点符号
        for i in range(len(truncated)-1, max_length//2, -1):
            if truncated[i] in [' ', '，', '。', '：', '-', '、']:
                return truncated[:i] + "..."
        
        return truncated + "..."
    
    def _format_summary_content(self, summary: str, max_length: int = 800) -> str:
        """格式化摘要内容以适应幻灯片"""
        if len(summary) <= max_length:
            return summary
        
        # 截断到最后一个完整句子
        truncated = summary[:max_length]
        last_period = truncated.rfind('。')
        if last_period > max_length * 0.7:  # 如果句号位置合理
            return truncated[:last_period + 1]
        else:
            return truncated + "..."
    
    def _format_summary_for_slide(self, summary: str) -> str:
        """专门为幻灯片格式化摘要内容"""
        # 16:9比例可以容纳更多内容，但仍需控制长度
        max_length = 800
        
        if len(summary) <= max_length:
            formatted = summary
        else:
            # 截断到最后一个完整句子
            truncated = summary[:max_length]
            last_period = truncated.rfind('。')
            if last_period > max_length * 0.7:
                formatted = truncated[:last_period + 1]
            else:
                formatted = truncated + "..."
        
        # 添加适当的段落分隔，提高可读性
        lines = formatted.split('。')
        formatted_lines = []
        
        for line in lines:
            line = line.strip()
            if line:
                if not line.endswith('。') and line != lines[-1]:
                    line += '。'
                formatted_lines.append(line)
        
        # 每3-4句为一段，适应16:9更宽的布局
        paragraphs = []
        current_paragraph = []
        
        for i, line in enumerate(formatted_lines):
            current_paragraph.append(line)
            if len(current_paragraph) >= 3 or i == len(formatted_lines) - 1:
                paragraphs.append(' '.join(current_paragraph))
                current_paragraph = []
        
        return '\n'.join(paragraphs)  # 减少段落间距
    
    def _add_url_textbox(self, slide, url: str):
        """在幻灯片底部添加原文链接文本框"""
        # 在16:9幻灯片底部添加文本框
        left = Inches(0.5)
        top = Inches(6.8)  # 适应16:9比例的底部位置
        width = Inches(12.0)  # 适应16:9比例的宽度
        height = Inches(0.6)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # 添加链接文本
        p = text_frame.paragraphs[0]
        p.text = f"🔗 原文链接: {url}"
        
        # 设置链接样式
        p.font.name = self.template_config.font_family
        p.font.size = Pt(self.template_config.url_font_size)
        p.font.color.rgb = RGBColor(0, 102, 204)  # 蓝色链接
        p.font.italic = True
        p.line_spacing = 1.0  # 减小行间距
    
    def save_presentation(self, filename: str) -> bool:
        """
        保存演示文稿到文件
        
        Args:
            filename: 文件名
            
        Returns:
            bool: 保存是否成功
        """
        try:
            if self.presentation is None:
                logging.error("没有可保存的演示文稿")
                return False
            
            self.presentation.save(filename)
            logging.info(f"PowerPoint演示文稿已保存: {filename}")
            return True
            
        except Exception as e:
            logging.error(f"保存PowerPoint演示文稿失败: {e}")
            return False

class SlideTemplateManager:
    """幻灯片模板管理器"""
    
    @staticmethod
    def load_default_template() -> TemplateConfig:
        """加载默认模板配置"""
        return TemplateConfig()
    
    @staticmethod
    def load_custom_template(config_path: str) -> Optional[TemplateConfig]:
        """
        从配置文件加载自定义模板
        
        Args:
            config_path: 配置文件路径
            
        Returns:
            TemplateConfig或None
        """
        try:
            import json
            with open(config_path, 'r', encoding='utf-8') as f:
                config_data = json.load(f)
            
            return TemplateConfig(**config_data)
        except Exception as e:
            logging.error(f"加载自定义模板失败: {e}")
            return None

class ContentFormatter:
    """内容格式化器"""
    
    @staticmethod
    def format_summary_content(summary: str) -> str:
        """格式化摘要内容"""
        # 简单的格式化：确保段落分隔
        lines = summary.split('\n')
        formatted_lines = []
        
        for line in lines:
            line = line.strip()
            if line:
                formatted_lines.append(line)
        
        return '\n\n'.join(formatted_lines)
    
    @staticmethod
    def create_bullet_points(content: str) -> str:
        """将内容转换为项目符号格式"""
        lines = content.split('\n')
        bullet_points = []
        
        for line in lines:
            line = line.strip()
            if line and not line.startswith('•'):
                bullet_points.append(f"• {line}")
            elif line:
                bullet_points.append(line)
        
        return '\n'.join(bullet_points)
    
    @staticmethod
    def truncate_content(content: str, max_length: int) -> str:
        """截断内容到指定长度"""
        if len(content) <= max_length:
            return content
        
        truncated = content[:max_length]
        last_space = truncated.rfind(' ')
        if last_space > max_length * 0.8:
            return truncated[:last_space] + "..."
        else:
            return truncated + "..."

def is_powerpoint_available() -> bool:
    """检查PowerPoint生成功能是否可用"""
    return PPTX_AVAILABLE